from rest_framework.decorators import api_view, permission_classes, parser_classes
from rest_framework.permissions import IsAuthenticated, AllowAny
from rest_framework.response import Response
from rest_framework import status
from django.contrib.auth.models import User
from django.contrib.auth import authenticate 
from rest_framework.authtoken.models import Token
from rest_framework.parsers import MultiPartParser, JSONParser
from django.conf import settings
import json
from openai import OpenAI
import PyPDF2
from pptx import Presentation
from .models import AIActivity, UserProfile, DailyGoal
from google.oauth2 import id_token
from google.auth.transport import requests as google_requests

# Initialize the AI Client
client = OpenAI(api_key=settings.OPENAI_API_KEY)
print("Loading AI Model... please wait...")

# User Registration
@api_view(['POST'])
def signup(request):
    try:
        email = request.data.get('email')
        password = request.data.get('password')

        if User.objects.filter(username=email).exists():
            return Response({'error': 'This email is already registered'}, status=status.HTTP_400_BAD_REQUEST)

        user = User.objects.create_user(username=email, email=email, password=password)
        return Response({'message': 'User created successfully!'}, status=status.HTTP_201_CREATED)

    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_400_BAD_REQUEST)

# User Authentication
@api_view(['POST'])
def login_view(request):
    email = request.data.get('email')
    password = request.data.get('password')
    user = authenticate(username=email, password=password)

    if user is not None:
        token, _ = Token.objects.get_or_create(user=user)
        return Response({
            'token': token.key,
            'message': 'Login Successful!',
            'is_admin': user.is_staff  
        }, status=status.HTTP_200_OK)
    else:
        return Response({'error': 'Invalid email or password'}, status=status.HTTP_401_UNAUTHORIZED)
    
# User Management
@api_view(['GET', 'DELETE'])
def manage_users(request, user_id=None):
    if request.method == 'GET':
        users = User.objects.all().values('id', 'username', 'email', 'is_staff', 'date_joined')
        return Response(users, status=status.HTTP_200_OK)
    
    if request.method == 'DELETE':
        try:
            user_to_delete = User.objects.get(id=user_id)
            user_to_delete.delete()
            return Response({'message': 'User deleted successfully'}, status=status.HTTP_200_OK)
        except User.DoesNotExist:
            return Response({'error': 'User not found'}, status=status.HTTP_404_NOT_FOUND)
        
@api_view(['POST'])
@permission_classes([AllowAny])
def google_login(request):
    token = request.data.get('token') 
    
    if not token:
        return Response({'error': 'No token provided'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        # Verify the token securely with Google
        idinfo = id_token.verify_oauth2_token(
            token, 
            google_requests.Request(), 
            '98535180949-jeip6rmteon6isj19a0hfleq2ncs7iv9.apps.googleusercontent.com' 
        )

        email = idinfo.get('email')
        first_name = idinfo.get('given_name', '')
        last_name = idinfo.get('family_name', '')
        picture = idinfo.get('picture', '')

        # Check if user exists, otherwise create a new account
        user, created = User.objects.get_or_create(
            email=email,
            defaults={
                'username': email.split('@')[0], 
                'first_name': first_name,
                'last_name': last_name
            }
        )

        # Generate the Token for React
        auth_token, _ = Token.objects.get_or_create(user=user)

        return Response({
            'token': auth_token.key,
            'email': user.email,
            'name': user.first_name,
            'avatar': picture
        }, status=status.HTTP_200_OK)

    except ValueError:
        return Response({'error': 'Invalid Google token'}, status=status.HTTP_401_UNAUTHORIZED)
        

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def summarize_text(request):
    text = request.data.get('text', '')
    file_name = request.data.get('file_name', 'Pasted Notes')
    
    # 1. Grab the new summary style choice from React (defaults to 'default' if missing)
    summary_type = request.data.get('type', 'default') 
    
    if not text:
        return Response({'error': 'No text provided'}, status=status.HTTP_400_BAD_REQUEST)

    # 2. Dynamically engineer the prompt based on what the user clicked
    if summary_type == 'short':
        system_instruction = "You are an expert study assistant. Summarize the following educational text in exactly 2 or 3 highly concise sentences. Get straight to the point."
    elif summary_type == 'bullets':
        system_instruction = "You are an expert study assistant. Summarize the following educational text using only clear, concise bullet points. Focus on key concepts and definitions. Do not write paragraphs."
    else:
        system_instruction = "You are an expert study assistant. Provide a well-structured, comprehensive summary of the following educational text, highlighting key concepts, definitions, and main ideas."

    try:
        # 3. Send the dynamic instruction to OpenAI
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_instruction},
                {"role": "user", "content": text}
            ]
        )
        
        summary_text = response.choices[0].message.content
        
        # 4. Save the activity to the database (Bonus: Adds the style to the title!)
        activity = AIActivity.objects.create(
            user=request.user,
            activity_type="Summary",
            title=f"Document Summary ({summary_type.capitalize()})", 
            file_name=file_name,
            content=summary_text 
        )

        return Response({'summary': summary_text, 'activity_id': activity.id}, status=status.HTTP_200_OK)

    except Exception as e:
        print("Summary Error:", str(e))
        return Response({'error': 'Failed to generate summary. Please try again.'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Generate Quiz
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def generate_quiz(request):
    text = request.data.get('text', '')
    file_name = request.data.get('file_name', 'Pasted Notes')
    
    if not text:
        return Response({'error': 'No text provided'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            response_format={ "type": "json_object" },
            messages=[
                {
                    "role": "system", 
                    "content": "You are a quiz generator. Generate exactly 10 multiple-choice questions based on the provided text. Return ONLY a valid JSON object matching this exact format: {\"quiz_array\": [{\"question\": \"Sample Question?\", \"options\": [\"A) First option\", \"B) Second option\", \"C) Third option\", \"D) Fourth option\"], \"answer\": \"A) First option\"}]}. The 'answer' field MUST exactly match the full text of the correct option, not just the letter. Do not use markdown blocks."
                },
                {"role": "user", "content": text}
            ]
        )
        
        raw_data = json.loads(response.choices[0].message.content)
        # Extract the exact array that React expects!
        quiz_array = raw_data.get('quiz_array', [])

        activity = AIActivity.objects.create(
            user=request.user,
            activity_type='Quiz',
            title='Interactive Quiz',
            file_name=file_name,
            content=json.dumps(quiz_array)
        )

        return Response({'quiz': quiz_array, 'activity_id': activity.id}, status=status.HTTP_200_OK)
        
    except Exception as e:
        print("Quiz Error:", str(e))
        return Response({'error': 'Failed to generate quiz. Please try again.'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)



@api_view(['POST'])
@permission_classes([IsAuthenticated])
def generate_flashcards(request):
    text = request.data.get('text', '')
    file_name = request.data.get('file_name', 'Pasted Notes')
    
    if not text:
        return Response({'error': 'No text provided'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            response_format={ "type": "json_object" },
           messages=[
                {
                    "role": "system", 
                    "content": "You are an expert study flashcard generator. Extract key terms from the provided text. Return ONLY a valid JSON object matching this exact format: {\"flashcard_array\": [{\"question\": \"Key Term\", \"answer\": \"Clear definition\"}]}. For the 'answer' field, provide a perfectly balanced definition that is strictly 1 to 2 sentences long. It must be easy to memorize, concise, but contain the most important core concept. Do not use markdown blocks."
                
                },
                {"role": "user", "content": text}
            ]
        )
        
        raw_data = json.loads(response.choices[0].message.content)
        #  Extract the exact array that React expects!
        flashcards_array = raw_data.get('flashcard_array', [])
        
        activity = AIActivity.objects.create(
            user=request.user,
            activity_type="Flashcards",
            title="Study Flashcards",
            file_name=file_name,
            content=json.dumps(flashcards_array)
        )

        return Response(flashcards_array, status=status.HTTP_200_OK)
        
    except Exception as e:
        print("Flashcard Error:", str(e))
        return Response({'error': 'Failed to generate flashcards. Please try again.'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    
# File Extraction Processing
@api_view(['POST'])
@parser_classes([MultiPartParser])
def extract_text_from_file(request):
    file_obj = request.FILES.get('file')
    
    if not file_obj:
        return Response({'error': 'No file uploaded'}, status=status.HTTP_400_BAD_REQUEST)

    filename = file_obj.name.lower()
    extracted_text = ""

    try:
        if filename.endswith('.txt'):
            extracted_text = file_obj.read().decode('utf-8')
        
        elif filename.endswith('.pdf'):
            pdf_reader = PyPDF2.PdfReader(file_obj)
            for page in pdf_reader.pages:
                extracted_text += page.extract_text() + "\n"
                
        elif filename.endswith('.pptx'):
            ppt = Presentation(file_obj)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        
        else:
            return Response({'error': 'Please upload a PDF, PPTX, or TXT file.'}, status=status.HTTP_400_BAD_REQUEST)

        return Response({'text': extracted_text.strip()}, status=status.HTTP_200_OK)

    except Exception as e:
        print("File Extraction Error:", str(e))
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Retrieve Dashboard Activity Data
@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_recent_activity(request):
    try:
        activities = AIActivity.objects.filter(user=request.user).order_by('-created_at')[:5]
        
        recent_list = []
        for act in activities:
            recent_list.append({
                "id": act.id,
                "type": act.activity_type,
                "title": act.title,
                "file": act.file_name,
                "time": act.created_at.strftime("%b %d, %Y - %H:%M"),
                "score": act.score, 
                "total": act.total_questions
            })
            
        total_quizzes = AIActivity.objects.filter(user=request.user, activity_type="Quiz").count()
        total_flashcards = AIActivity.objects.filter(user=request.user, activity_type="Flashcards").count()
        total_summaries = AIActivity.objects.filter(user=request.user, activity_type="Summary").count()
        total_files = AIActivity.objects.filter(user=request.user).exclude(file_name="Pasted Notes").values('file_name').distinct().count()

        profile_obj = UserProfile.objects.filter(user=request.user).first()
        avatar_url = None
        if profile_obj and profile_obj.profile_picture:
            avatar_url = request.build_absolute_uri(profile_obj.profile_picture.url)

        dashboard_data = {
            "recent_activity": recent_list,
            "progress_stats": {
                "quizzes_taken": total_quizzes,
                "flashcards_reviewed": total_flashcards,
                "summaries_generated": total_summaries,
                "files_uploaded": total_files
            },
            "username": request.user.first_name if request.user.first_name else request.user.username,
            "avatar": avatar_url
        }
            
        return Response(dashboard_data, status=status.HTTP_200_OK)
        
    except Exception as e: 
        print("Activity Fetch Error:", str(e))
        return Response({'error': 'Failed to fetch activity'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Save Quiz Metrics
@api_view(['POST'])
@permission_classes([IsAuthenticated])
def save_quiz_score(request):
    try:
        final_score = request.data.get('score')
        total_q = request.data.get('total')
        quiz_data = request.data.get('quiz_data')
        
        AIActivity.objects.create(
            user=request.user,
            activity_type="Quiz",
            title="Interactive Quiz",
            file_name="Completed Quiz",
            score=final_score,
            total_questions=total_q,
            content=json.dumps(quiz_data) if quiz_data else None
        )
        return Response({'message': 'Score saved successfully!'}, status=status.HTTP_200_OK)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

# Retrieve Specific Activity Record
@api_view(['GET'])
@permission_classes([IsAuthenticated])
def get_activity_detail(request, activity_id):
    try:
        activity = AIActivity.objects.get(id=activity_id, user=request.user)
        
        return Response({
            "id": activity.id,
            "type": activity.activity_type,
            "title": activity.title,
            "content": activity.content, 
            "time": activity.created_at.strftime("%b %d, %Y")
        }, status=status.HTTP_200_OK)
        
    except AIActivity.DoesNotExist:
        return Response({'error': 'Activity not found'}, status=status.HTTP_404_NOT_FOUND)

# User Profile Management
@api_view(['GET', 'PUT', 'DELETE'])
@permission_classes([IsAuthenticated])
@parser_classes([MultiPartParser, JSONParser])
def user_profile(request):
    profile_obj, created = UserProfile.objects.get_or_create(user=request.user)

    if request.method == 'GET':
        display_name = request.user.first_name if request.user.first_name else request.user.username.split('@')[0]
        
        avatar_url = None
        if profile_obj.profile_picture:
            avatar_url = request.build_absolute_uri(profile_obj.profile_picture.url)
            
        return Response({
            'email': request.user.email,
            'name': display_name,
            'avatar': avatar_url, 
            'date_joined': request.user.date_joined.strftime("%B %d, %Y")
        }, status=status.HTTP_200_OK)
        
    elif request.method == 'PUT':
        updated = False
        
        new_name = request.data.get('name')
        if new_name:
            request.user.first_name = new_name
            request.user.save()
            updated = True
            
        if 'profile_picture' in request.FILES:
            profile_obj.profile_picture = request.FILES['profile_picture']
            profile_obj.save()
            updated = True

        if updated:
            return Response({'message': 'Profile updated successfully!'}, status=status.HTTP_200_OK)
        else:
            return Response({'error': 'No data provided to update.'}, status=status.HTTP_400_BAD_REQUEST)
        
    elif request.method == 'DELETE':
        request.user.delete()
        return Response({'message': 'Account deleted successfully!'}, status=status.HTTP_200_OK)


@api_view(['GET', 'POST', 'PUT', 'DELETE'])
@permission_classes([IsAuthenticated])
def manage_goals(request, goal_id=None):
    if request.method == 'GET':
        # Get all goals for this user, newest first
        goals = DailyGoal.objects.filter(user=request.user).order_by('completed', '-created_at')
        data = [{"id": g.id, "text": g.text, "completed": g.completed} for g in goals]
        return Response(data, status=status.HTTP_200_OK)

    elif request.method == 'POST':
        # Create a new goal
        text = request.data.get('text')
        if not text:
            return Response({'error': 'Goal text is required'}, status=status.HTTP_400_BAD_REQUEST)
        goal = DailyGoal.objects.create(user=request.user, text=text)
        return Response({"id": goal.id, "text": goal.text, "completed": goal.completed}, status=status.HTTP_201_CREATED)

    elif request.method == 'PUT':
        # Toggle complete/incomplete
        try:
            goal = DailyGoal.objects.get(id=goal_id, user=request.user)
            goal.completed = not goal.completed 
            goal.save()
            return Response({"id": goal.id, "text": goal.text, "completed": goal.completed}, status=status.HTTP_200_OK)
        except DailyGoal.DoesNotExist:
            return Response({'error': 'Goal not found'}, status=status.HTTP_404_NOT_FOUND)

    elif request.method == 'DELETE':
        # Delete a goal entirely
        try:
            goal = DailyGoal.objects.get(id=goal_id, user=request.user)
            goal.delete()
            return Response({'message': 'Goal deleted'}, status=status.HTTP_200_OK)
        except DailyGoal.DoesNotExist:
            return Response({'error': 'Goal not found'}, status=status.HTTP_404_NOT_FOUND)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def quick_chat(request):
    user_message = request.data.get('message', '')
    
    if not user_message:
        return Response({'error': 'Message is required'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful, concise AI study assistant. Provide short, accurate answers to the student's questions."},
                {"role": "user", "content": user_message}
            ]
        )
        
        ai_reply = response.choices[0].message.content

        # Optional: Save it to their activity history so it shows up on the dashboard!
        AIActivity.objects.create(
            user=request.user,
            activity_type="Summary", # Using summary icon for chat
            title="Quick Q&A",
            file_name="General Chat",
            content=ai_reply
        )

        return Response({'reply': ai_reply}, status=status.HTTP_200_OK)
        
    except Exception as e:
        print("Chat Error:", str(e))
        return Response({'error': 'AI failed to respond.'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)